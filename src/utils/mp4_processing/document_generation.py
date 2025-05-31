# utils/document_generation.py
"""
Document Generation Utilities
Handles PowerPoint and document creation from processed media
"""

import re
import logging
from pathlib import Path
from datetime import datetime, timedelta
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt

from ..logger_setup import setup_logger

logger = setup_logger(name="config_manager", level=logging.INFO)


class DocumentGenerator:
    """Handles document generation including PowerPoint presentations."""
    
    def __init__(self):
        self.logger = logger
    
    def create_ppt_from_frames(self, frames_dir: Path, ppt_path: Path) -> bool:
        """
        Create PowerPoint presentation from extracted frames.
        
        Args:
            frames_dir: Directory containing frame images
            ppt_path: Output path for PowerPoint file
            
        Returns:
            bool: True if successful, False otherwise
        """
        self.logger.info(f"Starting PowerPoint creation from frames in {frames_dir}")
        images: List[Path] = sorted(frames_dir.glob("frame_*.png"))
        
        if not images:
            self.logger.warning("No frames found; skipping PPT generation")
            return False

        try:
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]  # completely blank
            slide_w = prs.slide_width
            slide_h = prs.slide_height

            ts_re = re.compile(r"frame_(\d{4})_\d+\.png")
            
            self.logger.info(f"Creating PowerPoint with {len(images)} frames")
            
            for i, img in enumerate(images):
                # Extract timestamp from filename
                m = ts_re.match(img.name)
                if m:
                    ts_raw = m.group(1)  # e.g. "0103"
                    mm, ss = ts_raw[:2], ts_raw[2:]
                    ts_label = f"{mm}:{ss}"
                else:
                    ts_label = "unknown"

                # Create slide
                slide = prs.slides.add_slide(blank_layout)

                # Add title textbox
                title_box = slide.shapes.add_textbox(
                    left=Inches(0.2),
                    top=Inches(0.1),
                    width=slide_w - Inches(0.4),
                    height=Inches(0.6)
                )
                p = title_box.text_frame.paragraphs[0]
                p.text = ts_label
                p.runs[0].font.size = Pt(24)
                p.runs[0].font.bold = True

                # Add picture - centered horizontally, flush to bottom
                pic = slide.shapes.add_picture(str(img), left=0, top=0, width=slide_w)
                pic.left = int((slide_w - pic.width) / 2)
                pic.top = int(slide_h - pic.height)
                
                if i % 10 == 0:
                    self.logger.info(f"Added slide {i+1} of {len(images)}")

            self.logger.info(f"Saving PowerPoint to {ppt_path}")
            prs.save(str(ppt_path))
            self.logger.info(f"PowerPoint created successfully: {ppt_path}")
            return True
            
        except Exception as e:
            self.logger.error(f"Error creating PowerPoint: {e}")
            return False
    
    def build_markdown_document(
        self,
        title: str,
        created: datetime,
        duration: timedelta,
        summary: str,
        transcript: str,
        additional_metadata: Optional[dict] = None
    ) -> str:
        """
        Build a markdown document with meeting information.
        
        Args:
            title: Document title
            created: Creation timestamp
            duration: Meeting duration
            summary: Meeting summary
            transcript: Meeting transcript
            additional_metadata: Optional additional metadata to include
            
        Returns:
            Formatted markdown string
        """
        self.logger.info("Building markdown document from metadata, summary, and transcript")
        
        # Format duration
        hours = int(duration.total_seconds() // 3600)
        minutes = int((duration.total_seconds() % 3600) // 60)
        duration_str = f"{hours}h {minutes}m" if hours > 0 else f"{minutes}m"
        
        # Build metadata section
        metadata_lines = [
            f"**Created:** {created:%Y-%m-%d %H:%M:%S}",
            f"**Duration:** {duration_str}"
        ]
        
        if additional_metadata:
            for key, value in additional_metadata.items():
                metadata_lines.append(f"**{key}:** {value}")
        
        metadata_section = "  \n".join(metadata_lines)
        
        # Build complete document
        md = f"""# {title}

        {metadata_section}

        ## Summary

        {summary}

        ## Transcript

        ```
        {transcript}
        ```"""
        
        self.logger.info("Markdown document built successfully")
        return md
    
    def create_output_structure(self, base_dir: str, subdirs: List[str]) -> dict:
        """
        Create output directory structure.
        
        Args:
            base_dir: Base output directory
            subdirs: List of subdirectories to create
            
        Returns:
            Dictionary mapping subdir names to Path objects
        """
        base_path = Path(base_dir)
        dirs = {}
        
        for subdir in subdirs:
            dir_path = base_path / subdir
            dir_path.mkdir(parents=True, exist_ok=True)
            dirs[subdir] = dir_path
            self.logger.info(f"Created directory: {dir_path}")
        
        return dirs